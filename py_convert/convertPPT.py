from pptx import Presentation

import argparse
import os
import re
import sys

import convertUtil

# TIMESTAMP for version information.
TIMESTAMP = "Version 2018-07-06"


# Process all text runs in Table elements
def processTable(shape, converter):
    conversionCount = 0
    rownum = 1
    for row in shape.table.rows:
        cellnum = 1
        for cell in row.cells:
            if cell.text_frame:
                for para in cell.text_frame.paragraphs:
                    if para.runs:
                        for run in para.runs:
                          fontObj = None
                          # TODO: Simplify this logic
                          if run.font:
                            fontObj = run.font
                          if fontObj and fontObj.name in converter.oldFonts:
                              tryResult = converter.convertText(run.text, None)
                              if tryResult != run.text:
                                conversionCount += 1
                                run.text = tryResult
                                fontObj.name = converter.unicodeFont
            cellnum += 1
        rownum += 1
    return conversionCount


# Process all text runs in Text Frame elements
def processTextFrame(shape, converter):
    conversionCount = 0
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
          fontObj = None
          if run.font:
            fontObj = run.font
          if fontObj and fontObj.name in converter.oldFonts:
            tryResult = converter.convertText(run.text, None)
            if tryResult != run.text:
              conversionCount += 1
              run.text = tryResult
            fontObj.name = converter.unicodeFont
    return conversionCount


# *****************************************************************
def processOnePresentation(path_to_presentation, output_dir,
                           converter):
  prs = Presentation(path_to_presentation)

  print( '%d slides found in %s' % (len(prs.slides), path_to_presentation))

  conversionCount = 0
  slideNum = 0
  for slide in prs.slides:
    newconversionCount = 0
    for shape in slide.shapes:
      if shape.has_table:
        newconversionCount += processTable(shape, converter)

      if shape.has_text_frame:
        newconversionCount += processTextFrame(shape, converter)
    print ('  Slide %d of %d has %d conversions' % (slideNum, len(prs.slides),
                                                    newconversionCount))

    slideNum += 1
    conversionCount += newconversionCount

  print ('  %d conversions applied to Unicode ' % conversionCount)

  if output_dir is not '':
    # String the directory tree to the file, substituting the output
    fileIn = os.path.split(path_to_presentation)[1]
    baseWOextension = os.path.splitext(fileIn)[0]
    new_path_to_presentation = os.path.join(
        output_dir, baseWOextension + '_unicode.pptx')
  else:
    baseWOextension = os.path.splitext(path_to_presentation)[0]
    new_path_to_presentation = os.path.join(
        output_dir, baseWOextension + '_unicode.pptx')
  prs.save(new_path_to_presentation)

  if conversionCount:
    print '  Output to file %s\n' % new_path_to_presentation
  else:
    print '  No new file written.\n'


def main(argv):
  args = convertUtil.parseArgs()

  path_to_presentation = args.filenames

  convertFileCount = 0
  for path in path_to_presentation:
    processOnePresentation(path, args.font, args.output_dir)
    convertFileCount += 1

  print ('\n\n%d files were processed' % convertFileCount)


if __name__ == "__main__":
  main(sys.argv)
