import sys

from pptx import Presentation

def processTable(shape):
    print 'SHAPE TABLE. %d rows' % len(shape.table.rows)
    rownum = 1
    for row in shape.table.rows:
        print '  row %d has %d cells' % (rownum, len(row.cells))
        cellnum = 1
        for cell in row.cells:
            if cell.text_frame:
                for para in cell.text_frame.paragraphs:
                    if para.runs:
                        print '  cell %d' % cellnum
                        for run in para.runs:
                            print('       %s' % run.text)
            cellnum += 1
        rownum += 1


if len(sys.argv) > 1:
    path_to_presentation = sys.argv[1]
else:
    path_to_presentation = "modified_mog.pptx"

print 'Data from file %s' % path_to_presentation

prs = Presentation(path_to_presentation)

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

print( '%d slides in %s' % (len(prs.slides), path_to_presentation))
# Restrict to first slide for now.
slideNum = 1
for slide in prs.slides:
    print ('$$$ %d shapes in slide' % slideNum)
    for shape in slide.shapes:

        if shape.has_table:
            processTable(shape)

        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
                print('  %s' % run.text)
    slideNum += 1

print('%d text runs found' % len(text_runs))
