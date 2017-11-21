from pptx import Presentation

import argparse
import os
import re
import sys

import osageConversion
import convertUtil

# Rule for detecting Latin text or Old Osage font.
# Some Old Osage text is in Latin CAPS, but with lower case a, e, and o.
# Limitation: [aeo] right after capital Latin is converted.
latinOsagePattern2 = r'[\^A-Z\[\]][A-Zaeo; \[\]\^\\\'\/\._`,!]+'

# To avoid converting English words
notOsageLatinLower = re.compile(r'[b-df-np-z]')

# This identifies traditional Osage private use characters
traditionalOsageCharacters = ur'([\uf020-\uf05e]+)'

# Font names:
OfficialOsageFont = 'Official Osage Language'


def replFunc(matchObj):
  if matchObj.group(0):
     # Avoid converting strings with [aeo] after capital, but in English
    if notOsageLatinLower.search(matchObj.group(0)):
      return matchObj.group(0)
    else:
      return osageConversion.oldOsageToUnicode(matchObj.group(0))

# Process all text runs in Table elements
def processTable(shape, outFont):
    conversionCount = 0
    rownum = 1
    for row in shape.table.rows:
        cellnum = 1
        for cell in row.cells:
            if cell.text_frame:
                for para in cell.text_frame.paragraphs:
                    if para.runs:
                        for run in para.runs:
                          fontObj = None
                          if run.font:
                            fontObj = run.font
                          if fontObj and fontObj.name == OfficialOsageFont:
                            if not notOsageLatinLower.search(run.text):
                              tryResult = re.subn(latinOsagePattern2, replFunc, run.text)
                              if tryResult[1]:
                                conversionCount += 1
                                run.text = tryResult[0]
                                fontObj.name = outFont
            cellnum += 1
        rownum += 1
    return conversionCount


# Process all text runs in Text Frame elements
def processsTextFrame(shape, outFont):
    conversionCount = 0
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
          fontObj = None
          if run.font:
            fontObj = run.font
          if fontObj and fontObj.name == OfficialOsageFont:
            tryResult = re.subn(latinOsagePattern2, replFunc, run.text)
            if tryResult[1]:
              conversionCount += 1
              run.text = tryResult[0]
              fontObj.name = outFont
    return conversionCount


def processOnePresentation(path_to_presentation, outputFont, output_dir=''):
  prs = Presentation(path_to_presentation)

  print( '%d slides found in %s' % (len(prs.slides), path_to_presentation))

  conversionCount = 0
  slideNum = 0
  for slide in prs.slides:
    newconversionCount = 0
    for shape in slide.shapes:
      if shape.has_table:
        newconversionCount += processTable(shape, outputFont)

      if shape.has_text_frame:
        newconversionCount += processsTextFrame(shape, outputFont)
    print ('  Slide %d of %d has %d conversions' % (slideNum, len(prs.slides),
                                                    newconversionCount))

    slideNum += 1
    conversionCount += newconversionCount

  print ('  %d conversions applied to Osage Unicode ' % conversionCount)

  if output_dir is not '':
    # String the directory tree to the file, substituting the output
    fileIn = os.path.split(path_to_presentation)[1]
    baseWOextension = os.path.splitext(fileIn)[0]
    new_path_to_presentation = os.path.join(
        output_dir, baseWOextension + '_unicode.pptx')
  else:
    baseWOextension = os.path.splitext(path_to_presentation)[0]
    new_path_to_presentation = os.path.join(
        output_dir, baseWOextension + '_unicode.pptx')
  prs.save(new_path_to_presentation)

  if conversionCount:
    print '  Output to file %s\n' % new_path_to_presentation
  else:
    print '  No new file written.\n'


def main(argv):
  args = convertUtil.parseArgs()

  path_to_presentation = args.filenames

  convertFileCount = 0
  for path in path_to_presentation:
    processOnePresentation(path, args.font, args.output_dir)
    convertFileCount += 1

  print ('\n\n%d files were processed' % convertFileCount)


if __name__ == "__main__":
  main(sys.argv)
